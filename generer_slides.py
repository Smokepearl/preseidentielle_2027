# Génère le support de présentation (PowerPoint) pour l'oral
# Projet fil rouge - Présidentielle 2027 (LFI)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette (cohérente avec l'app) ────────────────────────────────────────
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_CLAIR = RGBColor(0x8B, 0x5C, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
CORAIL = RGBColor(0xFF, 0x4D, 0x6D)
SOMBRE = RGBColor(0x1E, 0x1B, 0x36)
GRIS = RGBColor(0x55, 0x55, 0x66)
GRIS_CLAIR = RGBColor(0xF4, 0xF2, 0xFB)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)

LARGEUR = Inches(13.333)
HAUTEUR = Inches(7.5)

prs = Presentation()
prs.slide_width = LARGEUR
prs.slide_height = HAUTEUR
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, couleur, ligne=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = couleur
    if ligne is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = ligne
    sh.shadow.inherit = False
    return sh


def texte(slide, x, y, w, h, contenu, taille=18, couleur=SOMBRE, gras=False,
          align=PP_ALIGN.LEFT, italique=False, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = contenu
    run.font.size = Pt(taille)
    run.font.color.rgb = couleur
    run.font.bold = gras
    run.font.italic = italique
    run.font.name = "Calibri"
    return tb


def puces(slide, x, y, w, h, items, taille=16, couleur=SOMBRE, interligne=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, niveau) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = niveau
        p.line_spacing = interligne
        p.space_after = Pt(6)
        run = p.add_run()
        prefixe = "•  " if niveau == 0 else "–  "
        run.text = prefixe + txt
        run.font.size = Pt(taille if niveau == 0 else taille - 2)
        run.font.color.rgb = couleur if niveau == 0 else GRIS
        run.font.bold = (niveau == 0)
        run.font.name = "Calibri"
    return tb


def entete_contenu(slide, titre, numero):
    # bande latérale + titre
    rect(slide, 0, 0, Inches(0.25), HAUTEUR, VIOLET)
    rect(slide, Inches(0.25), 0, LARGEUR - Inches(0.25), Inches(1.1), GRIS_CLAIR)
    texte(slide, Inches(0.55), Inches(0.18), Inches(11), Inches(0.8), titre,
          taille=28, couleur=VIOLET, gras=True, anchor=MSO_ANCHOR.MIDDLE)
    # pastille numéro
    pa = slide.shapes.add_shape(MSO_SHAPE.OVAL, LARGEUR - Inches(1.05),
                                Inches(0.32), Inches(0.5), Inches(0.5))
    pa.fill.solid(); pa.fill.fore_color.rgb = VIOLET; pa.line.fill.background()
    pa.shadow.inherit = False
    pa.text_frame.text = str(numero)
    pa.text_frame.paragraphs[0].runs[0].font.color.rgb = BLANC
    pa.text_frame.paragraphs[0].runs[0].font.bold = True
    pa.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    pa.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def slide_contenu(titre, numero):
    s = prs.slides.add_slide(BLANK)
    entete_contenu(s, titre, numero)
    return s


def slide_section(titre, sous_titre=""):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, LARGEUR, HAUTEUR, VIOLET)
    rect(s, 0, Inches(3.0), LARGEUR, Inches(0.06), CYAN)
    texte(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1), titre,
          taille=40, couleur=BLANC, gras=True, align=PP_ALIGN.CENTER)
    if sous_titre:
        texte(s, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8), sous_titre,
              taille=18, couleur=RGBColor(0xD6, 0xCB, 0xFF), align=PP_ALIGN.CENTER)
    return s


def tableau(slide, x, y, w, h, donnees, largeurs=None, taille=12,
            couleur_entete=VIOLET):
    n_lignes = len(donnees)
    n_cols = len(donnees[0])
    gtab = slide.shapes.add_table(n_lignes, n_cols, x, y, w, h).table
    if largeurs:
        for j, lw in enumerate(largeurs):
            gtab.columns[j].width = lw
    for i, ligne in enumerate(donnees):
        for j, val in enumerate(ligne):
            cell = gtab.cell(i, j)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(taille)
            para.runs[0].font.name = "Calibri"
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.margin_left = Pt(6); cell.margin_right = Pt(6)
            if i == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = couleur_entete
                para.runs[0].font.color.rgb = BLANC
                para.runs[0].font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLANC if i % 2 else GRIS_CLAIR
                para.runs[0].font.color.rgb = SOMBRE
    return gtab


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, LARGEUR, HAUTEUR, SOMBRE)
rect(s, 0, 0, LARGEUR, Inches(0.18), VIOLET)
rect(s, 0, HAUTEUR - Inches(0.18), LARGEUR, Inches(0.18), CYAN)
texte(s, Inches(1), Inches(1.2), Inches(11.3), Inches(0.6),
      "PROJET FIL ROUGE — B3 DATA & IA", taille=18, couleur=CYAN, gras=True,
      align=PP_ALIGN.CENTER)
texte(s, Inches(1), Inches(2.3), Inches(11.3), Inches(1.4),
      "Présidentielle 2027", taille=54, couleur=BLANC, gras=True, align=PP_ALIGN.CENTER)
texte(s, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
      "Analyse électorale & data storytelling", taille=24,
      couleur=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)
texte(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6),
      "Mission pour La France Insoumise", taille=18, couleur=CORAIL,
      align=PP_ALIGN.CENTER, italique=True)
texte(s, Inches(1), Inches(5.6), Inches(11.3), Inches(0.9),
      "H. El Idrissi  &  [Nom du binôme]\nSoutenance orale — [date]", taille=16,
      couleur=RGBColor(0xAA, 0xAA, 0xC0), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("Sommaire", 1)
gauche = [
    ("1.  Contexte & objectifs", 0),
    ("2.  Les données", 0),
    ("3.  Méthodologie du projet", 0),
    ("4.  Préparation des données", 0),
    ("5.  Analyse exploratoire", 0),
]
droite = [
    ("6.  Modélisation & simulateur", 0),
    ("7.  L'application interactive", 0),
    ("8.  Résultats & conclusions", 0),
    ("9.  Organisation & répartition", 0),
    ("10. Bilan & perspectives", 0),
]
puces(s, Inches(1.0), Inches(1.6), Inches(5.7), Inches(5), gauche, taille=20, interligne=1.5)
puces(s, Inches(7.0), Inches(1.6), Inches(5.7), Inches(5), droite, taille=20, interligne=1.5)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTE & MISSION
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("Contexte & mission", 2)
texte(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.9),
      "Sujet 1 : « Présidentielle 2027 — Analyse électorale et prédictions ». "
      "Mission freelance pour un cabinet de conseil politique, ici mandaté par "
      "La France Insoumise (LFI).", taille=17, couleur=GRIS)
# encadré problématique
rect(s, Inches(0.7), Inches(2.5), Inches(11.9), Inches(1.1), GRIS_CLAIR)
rect(s, Inches(0.7), Inches(2.5), Inches(0.12), Inches(1.1), CORAIL)
texte(s, Inches(1.0), Inches(2.62), Inches(11.4), Inches(0.9),
      "Problématique : Mélenchon (ou un candidat LFI) peut-il se qualifier "
      "au 2ᵉ tour de la présidentielle 2027 ?", taille=19, couleur=SOMBRE, gras=True,
      anchor=MSO_ANCHOR.MIDDLE)
puces(s, Inches(0.7), Inches(3.9), Inches(11.9), Inches(3), [
    ("Analyser les présidentielles passées (2012, 2017, 2022)", 0),
    ("Identifier les dynamiques régionales et socio-démographiques du vote", 0),
    ("Construire un modèle / simulateur pour explorer des scénarios 2027", 0),
    ("Livrer une application interactive de data storytelling (Streamlit)", 0),
], taille=18, interligne=1.3)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — LES DONNÉES
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("Les données", 3)
texte(s, Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.5),
      "Sources officielles, organisées en données brutes puis nettoyées.",
      taille=16, couleur=GRIS, italique=True)
tableau(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(2.6), [
    ["Scrutin", "Source", "Niveau", "Usage"],
    ["Présidentielle 2012", "Ministère Intérieur / data.gouv", "Régions", "Analyse"],
    ["Présidentielle 2017", "Ministère Intérieur / data.gouv", "Bureaux → régions", "Analyse"],
    ["Présidentielle 2022", "Ministère Intérieur / data.gouv", "Régions", "Analyse + base modèle"],
    ["Européennes 2024", "Ministère Intérieur / data.gouv", "Régions", "Signal de mobilisation"],
], largeurs=[Inches(3.0), Inches(4.2), Inches(2.5), Inches(2.2)], taille=13)
puces(s, Inches(0.7), Inches(4.8), Inches(11.9), Inches(2.3), [
    ("Données sociales INSEE par région : taux de chômage + niveau de vie (croisées avec le vote)", 0),
    ("Sociologie du vote : Ipsos / Sopra Steria (sortie des urnes 2022)", 0),
    ("Sondages 2027 : Odoxa (mai 2026)", 0),
    ("Fond de carte : GeoJSON des régions (france-geojson)", 0),
], taille=16, interligne=1.2)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MÉTHODOLOGIE (PIPELINE)
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("Méthodologie du projet", 4)
etapes = [
    ("1", "Collecte", "data.gouv, INSEE,\nsondages"),
    ("2", "Nettoyage", "brut → clean,\nstandardisation"),
    ("3", "Analyse", "exploration,\nvisualisations"),
    ("4", "Modélisation", "régression,\nscénarios 2027"),
    ("5", "Application", "Streamlit,\nstorytelling"),
]
x0 = Inches(0.7)
larg = Inches(2.25)
ecart = Inches(2.42)
for i, (num, titre, desc) in enumerate(etapes):
    x = Emu(int(x0) + i * int(ecart))
    rect(s, x, Inches(2.4), larg, Inches(2.2), VIOLET if i % 2 == 0 else VIOLET_CLAIR)
    texte(s, x, Inches(2.55), larg, Inches(0.7), num, taille=30, couleur=BLANC,
          gras=True, align=PP_ALIGN.CENTER)
    texte(s, x, Inches(3.25), larg, Inches(0.6), titre, taille=17, couleur=BLANC,
          gras=True, align=PP_ALIGN.CENTER)
    texte(s, x, Inches(3.8), larg, Inches(0.8), desc, taille=12,
          couleur=RGBColor(0xE5, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
    if i < len(etapes) - 1:
        fleche = slide_contenu  # noop
texte(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.4),
      "Un pipeline data classique, du fichier brut jusqu'à l'application livrée. "
      "Chaque étape est documentée dans un notebook dédié et versionnée dans le projet.",
      taille=16, couleur=GRIS, italique=True)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PRÉPARATION DES DONNÉES
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("Préparation des données", 5)
puces(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(2.6), [
    ("Organisation en deux dossiers : data_brut/ et data_clean/", 0),
    ("Un notebook de construction par scrutin (2012, 2017, européennes 2024)", 0),
    ("Standardisation : mêmes colonnes pour comparer les années", 0),
], taille=18, interligne=1.25)
# encadré défis
rect(s, Inches(0.7), Inches(3.7), Inches(11.9), Inches(2.9), GRIS_CLAIR)
rect(s, Inches(0.7), Inches(3.7), Inches(0.12), Inches(2.9), CYAN)
texte(s, Inches(1.0), Inches(3.85), Inches(11), Inches(0.5),
      "Difficultés techniques traitées :", taille=17, couleur=VIOLET, gras=True)
puces(s, Inches(1.0), Inches(4.4), Inches(11.3), Inches(2.1), [
    ("2012 : 22 anciennes régions → agrégation vers les 13 régions actuelles (réforme 2015)", 0),
    ("Format ministère « large » : 11-38 candidats/listes en colonnes → parsing par position", 0),
    ("Fichier mal nommé détecté : « 2017_bureaux » contenait en fait les données 2022", 0),
    ("Réseau de l'école (SSL) bloquant : fond de carte téléchargé puis stocké en local", 0),
], taille=14, interligne=1.15)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ANALYSE EXPLORATOIRE
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("Analyse exploratoire", 6)
# 3 chiffres clés
kpis = [("+10,9 pts", "progression de Mélenchon\n2012 → 2022"),
        ("3ᵉ", "place en 2022, à 1,2 pt\ndu 2ᵉ tour"),
        ("54 %", "de l'électorat LFI\na moins de 35 ans")]
for i, (val, lbl) in enumerate(kpis):
    x = Emu(int(Inches(0.7)) + i * int(Inches(4.05)))
    rect(s, x, Inches(1.5), Inches(3.7), Inches(1.7), GRIS_CLAIR)
    texte(s, x, Inches(1.6), Inches(3.7), Inches(0.8), val, taille=32, couleur=CORAIL,
          gras=True, align=PP_ALIGN.CENTER)
    texte(s, x, Inches(2.45), Inches(3.7), Inches(0.7), lbl, taille=13, couleur=GRIS,
          align=PP_ALIGN.CENTER)
puces(s, Inches(0.7), Inches(3.6), Inches(11.9), Inches(3), [
    ("Progression continue dans toutes les régions, surtout en Île-de-France et Occitanie", 0),
    ("Électorat jeune, urbain et populaire (fort chez les chômeurs : 30 %)", 0),
    ("La participation recule (-5,8 pts entre 2012 et 2022) : un réservoir de voix abstentionnistes", 0),
    ("Croisement INSEE : LFI plus fort dans les régions au niveau de vie élevé (r=0,70) "
     "= « effet métropole » (et non lié au chômage régional)", 0),
], taille=16, interligne=1.25)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MODÉLISATION & SIMULATEUR
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("Modélisation & simulateur", 7)
puces(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(2.4), [
    ("Régression linéaire sur la tendance 2012-2017-2022 (projection prudente)", 0),
    ("Simulateur de scénarios : union à gauche, participation, mobilisation des jeunes, "
     "score du RN, type de candidat", 0),
    ("Signal de mobilisation 2024 : l'écart présidentielle/européennes (~12 pts) calibre "
     "l'effet de la mobilisation, sans mélanger les types de scrutin", 0),
], taille=17, interligne=1.3)
tableau(s, Inches(0.7), Inches(4.2), Inches(11.9), Inches(2.2), [
    ["Scénario", "Conditions", "Score estimé", "2ᵉ tour ?"],
    ["Pessimiste", "LFI seule, faible mobilisation", "~18 %", "Incertain"],
    ["Réaliste", "Union légère + mobilisation normale", "~22-24 %", "Probable"],
    ["Optimiste", "Union totale + forte mobilisation", "~26-28 %", "Très probable"],
], largeurs=[Inches(2.2), Inches(5.3), Inches(2.2), Inches(2.2)], taille=14)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — L'APPLICATION
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("L'application interactive", 8)
texte(s, Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.6),
      "Application Streamlit en récit guidé — 7 chapitres, design « dashboard » moderne.",
      taille=17, couleur=GRIS, italique=True)
chapitres = ["🎯 Contexte", "📈 L'ascension", "🗺️ La carte", "👥 Qui vote LFI",
             "🇪🇺 Signal 2024", "🔮 Simuler 2027", "✅ Le verdict"]
for i, ch in enumerate(chapitres):
    col = i % 4
    lig = i // 4
    x = Emu(int(Inches(0.7)) + col * int(Inches(3.0)))
    y = Emu(int(Inches(2.1)) + lig * int(Inches(1.0)))
    rect(s, x, y, Inches(2.8), Inches(0.8), VIOLET_CLAIR)
    texte(s, x, y, Inches(2.8), Inches(0.8), ch, taille=15, couleur=BLANC,
          gras=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
puces(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(2.3), [
    ("Carte interactive, filtres par année et indicateur", 0),
    ("Simulateur 2027 avec paramètres modifiables en direct", 0),
    ("Déploiement local + documentation (README, DATASETS)", 0),
    ("Démonstration en direct pendant la soutenance", 0),
], taille=16, interligne=1.2)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — RÉSULTATS & CONCLUSIONS
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("Résultats & conclusions", 9)
texte(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.6),
      "Les 3 leviers décisifs pour LFI en 2027 :", taille=18, couleur=SOMBRE, gras=True)
leviers = [("🎯 Mobiliser", "Abstention à 29 % dans\nles quartiers populaires.\n+3 à +5 pts"),
           ("🤝 Unir la gauche", "~9 pts de gauche hors LFI\nen 2022 à reconquérir.\n+3 à +6 pts"),
           ("👥 Capter les jeunes", "30 % des -35 ans votent\nLFI mais s'abstiennent.\n+1 à +2 pts")]
for i, (titre, desc) in enumerate(leviers):
    x = Emu(int(Inches(0.7)) + i * int(Inches(4.05)))
    rect(s, x, Inches(2.1), Inches(3.7), Inches(2.0), GRIS_CLAIR)
    rect(s, x, Inches(2.1), Inches(3.7), Inches(0.1), CORAIL)
    texte(s, x, Inches(2.3), Inches(3.7), Inches(0.6), titre, taille=18, couleur=VIOLET,
          gras=True, align=PP_ALIGN.CENTER)
    texte(s, x, Inches(2.95), Inches(3.7), Inches(1.1), desc, taille=14, couleur=GRIS,
          align=PP_ALIGN.CENTER)
rect(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(1.6), VIOLET)
texte(s, Inches(1.0), Inches(4.75), Inches(11.4), Inches(1.3),
      "En résumé : avec +10,9 pts en 10 ans, LFI a les moyens d'atteindre le 2ᵉ tour. "
      "La clé tient en deux mots — union et mobilisation. "
      "Et les sondages sous-estiment historiquement LFI (9 % annoncés en 2022, 22 % réels).",
      taille=17, couleur=BLANC, gras=True, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CALENDRIER DU PROJET
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("Calendrier du projet", 10)
texte(s, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.5),
      "Déroulé sur 6 semaines (à adapter à vos dates réelles).",
      taille=15, couleur=GRIS, italique=True)
tableau(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(4.6), [
    ["Semaine", "Phase", "Livrable principal"],
    ["S1", "Cadrage : choix du sujet, recherche des sources", "Plan de projet"],
    ["S2", "Collecte & nettoyage des données", "Fichiers clean (2012/17/22/24)"],
    ["S3", "Analyse exploratoire & visualisations", "Notebook d'analyse"],
    ["S4", "Modélisation & scénarios 2027", "Notebook de modélisation"],
    ["S5", "Développement de l'application Streamlit", "App interactive"],
    ["S6", "Design, documentation & préparation orale", "App finale + slides"],
], largeurs=[Inches(1.4), Inches(6.5), Inches(4.0)], taille=14)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — RÉPARTITION DES TÂCHES (BINÔME)
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("Répartition des tâches (binôme)", 11)
texte(s, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.5),
      "Travail en binôme, avec mise en commun régulière. Répartition indicative — à ajuster.",
      taille=15, couleur=GRIS, italique=True)
tableau(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.4), [
    ["Phase", "H. El Idrissi", "[Binôme]"],
    ["Cadrage / recherche données", "Pilote", "Appui"],
    ["Collecte & nettoyage", "Construction datasets", "Vérification / contrôle qualité"],
    ["Analyse exploratoire", "Croisements socio", "Visualisations"],
    ["Modélisation & scénarios", "Modèle + simulateur", "Calibrage / hypothèses"],
    ["Application Streamlit", "Développement", "Design & storytelling"],
    ["Documentation & oral", "README / datasets", "Slides & démo"],
], largeurs=[Inches(4.3), Inches(3.8), Inches(3.8)], taille=13)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — BILAN & COMPÉTENCES
# ══════════════════════════════════════════════════════════════════════════
s = slide_contenu("Bilan & perspectives", 12)
texte(s, Inches(0.7), Inches(1.35), Inches(5.6), Inches(0.5),
      "Compétences mobilisées", taille=18, couleur=VIOLET, gras=True)
puces(s, Inches(0.7), Inches(1.9), Inches(5.6), Inches(3.5), [
    ("Collecte & nettoyage (pandas)", 0),
    ("Visualisation (Plotly, Matplotlib)", 0),
    ("Machine learning (scikit-learn)", 0),
    ("Développement web (Streamlit)", 0),
    ("Data storytelling & restitution", 0),
    ("Travail en binôme & versioning", 0),
], taille=16, interligne=1.25)
texte(s, Inches(6.7), Inches(1.35), Inches(5.9), Inches(0.5),
      "Limites & pistes", taille=18, couleur=VIOLET, gras=True)
puces(s, Inches(6.7), Inches(1.9), Inches(5.9), Inches(3.5), [
    ("Modèle simple (3 points) : tendance, pas une prédiction", 0),
    ("Croiser davantage avec données INSEE par commune", 0),
    ("Ajouter les législatives 2024 comme signal", 0),
    ("Déploiement en ligne (Streamlit Cloud)", 0),
], taille=16, interligne=1.25)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — MERCI
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, LARGEUR, HAUTEUR, SOMBRE)
rect(s, 0, Inches(3.5), LARGEUR, Inches(0.06), CYAN)
texte(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1), "Merci de votre attention",
      taille=40, couleur=BLANC, gras=True, align=PP_ALIGN.CENTER)
texte(s, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8), "Questions & démonstration",
      taille=22, couleur=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)
texte(s, Inches(1), Inches(5.0), Inches(11.3), Inches(0.6),
      "Présidentielle 2027 — Analyse électorale · B3 DATA & IA", taille=14,
      couleur=RGBColor(0xAA, 0xAA, 0xC0), align=PP_ALIGN.CENTER)


prs.save("Presentation_orale_LFI_2027.pptx")
print("OK - Presentation_orale_LFI_2027.pptx genere :", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
